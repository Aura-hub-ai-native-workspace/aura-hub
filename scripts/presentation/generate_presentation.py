from __future__ import annotations

"""Generate the AURA Hub architecture presentation.

The deck is intentionally built from editable PowerPoint shapes rather than
screenshots or raster assets.  The content mirrors the current working tree:
implemented systems are called out as live, while partial integrations and
roadmap work are explicitly separated.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("presentation.pptx")

W = 13.333
H = 7.5


def C(value: str) -> RGBColor:
    value = value.replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


BG = C("080C14")
DARK_ACCENT = C("1D2A40")
SURFACE = C("111827")
SURFACE_2 = C("151F31")
SURFACE_3 = C("1A263A")
LINE = C("27344B")
TEXT = C("F1F5FF")
MUTED = C("91A0BA")
FAINT = C("60708D")
BLUE = C("5E83FF")
BLUE_DARK = C("3159D5")
CYAN = C("4BD8F7")
MINT = C("47DCA2")
AMBER = C("FFBF5B")
RED = C("FF6C67")
PURPLE = C("AE8BFF")

FONT = "Aptos"
DISPLAY = "Aptos Display"


def set_fill(shape, color: RGBColor | None, transparency: int = 0) -> None:
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.fill.transparency = transparency


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor | None = SURFACE,
    line: RGBColor | None = LINE,
    line_width: float = 1.0,
    transparency: int = 0,
    radius: bool = False,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else shape_type
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill, transparency)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def rect(slide, x, y, w, h, fill=SURFACE, line=LINE, radius=True, **kwargs):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line, radius=radius, **kwargs)


def circle(slide, x, y, d, fill=None, line=LINE, line_width=1.0, transparency=0):
    return add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill, line, line_width, transparency)


def line(slide, x1, y1, x2, y2, color=LINE, width=1.2, dash=None):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    if dash:
        connector.line.dash_style = dash
    return connector


def arrow(slide, x1, y1, x2, y2, color=BLUE, width=1.5, head=0.10):
    line(slide, x1, y1, x2, y2, color, width)
    if abs(y2 - y1) >= abs(x2 - x1):
        tri = add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x2 - head / 2, y2 - head / 2, head, head, color, None)
        tri.rotation = 180 if y2 > y1 else 0
    else:
        tri = add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x2 - head / 2, y2 - head / 2, head, head, color, None)
        tri.rotation = 90 if x2 > x1 else 270
    return tri


def text(
    slide,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 12,
    color: RGBColor = TEXT,
    bold: bool = False,
    font: str = FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0,
    italic: bool = False,
    caps: bool = False,
    line_spacing: float = 1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    lines = value.split("\n")
    for index, raw in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        run = paragraph.add_run()
        run.text = raw.upper() if caps else raw
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def pill(slide, value: str, x: float, y: float, w: float, color=BLUE, fill=None, size=8.5):
    fill = fill or color
    rect(slide, x, y, w, 0.25, fill=fill, line=None, radius=True)
    text(slide, value, x + 0.07, y + 0.035, w - 0.14, 0.16, size, BG, True, align=PP_ALIGN.CENTER, caps=True)


def status_pill(slide, value: str, x: float, y: float, w: float, tone=MINT):
    rect(slide, x, y, w, 0.27, fill=tone, line=None, radius=True)
    circle(slide, x + 0.08, y + 0.085, 0.09, BG, None)
    text(slide, value, x + 0.22, y + 0.045, w - 0.27, 0.16, 8.5, BG, True, caps=True)


def footer(slide, number: int):
    line(slide, 0.76, 7.08, 12.56, 7.08, C("1B2639"), 0.7)
    text(slide, "AURA HUB  /  ENGINEERING OPERATING ENVIRONMENT  /  REPO-VERIFIED", 0.78, 7.17, 7.6, 0.16, 7.3, FAINT, True, caps=True)
    text(slide, f"{number:02d}", 12.08, 7.14, 0.48, 0.2, 8.5, MUTED, True, align=PP_ALIGN.RIGHT)


def base(slide, number: int, eyebrow: str, title: str, subtitle: str, tag: str | None = None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    # restrained architectural grid
    for x in (0.8, 3.3, 5.8, 8.3, 10.8, 12.55):
        line(slide, x, 0.36, x, 6.92, C("0E1726"), 0.55)
    line(slide, 0.76, 0.36, 12.56, 0.36, C("142137"), 0.6)
    text(slide, eyebrow, 0.78, 0.51, 7.8, 0.18, 8.2, CYAN, True, caps=True)
    text(slide, title, 0.78, 0.76, 10.7, 0.78, 25.5, TEXT, True, font=DISPLAY)
    text(slide, subtitle, 0.80, 1.70, 10.9, 0.22, 11.2, MUTED)
    if tag:
        pill(slide, tag, 10.95, 0.52, 1.55, MINT)
    footer(slide, number)


def section_label(slide, value, x, y, w=2.0, color=FAINT):
    text(slide, value, x, y, w, 0.16, 7.8, color, True, caps=True)


def draw_aura_mark(slide, cx, cy, scale=1.0, accent=BLUE):
    # A geometric mark made from rings and a central aperture.
    circle(slide, cx - 0.72 * scale, cy - 0.72 * scale, 1.44 * scale, None, C("28436D"), 1.0)
    circle(slide, cx - 0.50 * scale, cy - 0.50 * scale, 1.00 * scale, None, accent, 1.8)
    circle(slide, cx - 0.20 * scale, cy - 0.20 * scale, 0.40 * scale, accent, None)
    line(slide, cx - 0.95 * scale, cy, cx - 0.52 * scale, cy, CYAN, 1.2)
    line(slide, cx + 0.52 * scale, cy, cx + 0.95 * scale, cy, CYAN, 1.2)
    line(slide, cx, cy - 0.95 * scale, cx, cy - 0.52 * scale, CYAN, 1.2)
    line(slide, cx, cy + 0.52 * scale, cx, cy + 0.95 * scale, CYAN, 1.2)


def draw_code_icon(slide, x, y, s=0.22, color=BLUE):
    line(slide, x + s * 0.34, y + s * 0.07, x + s * 0.09, y + s * 0.50, color, 1.2)
    line(slide, x + s * 0.09, y + s * 0.50, x + s * 0.34, y + s * 0.93, color, 1.2)
    line(slide, x + s * 0.66, y + s * 0.07, x + s * 0.91, y + s * 0.50, color, 1.2)
    line(slide, x + s * 0.91, y + s * 0.50, x + s * 0.66, y + s * 0.93, color, 1.2)
    line(slide, x + s * 0.57, y + s * 0.03, x + s * 0.43, y + s * 0.97, CYAN, 1.0)


def draw_graph_icon(slide, x, y, s=0.28, color=CYAN):
    line(slide, x + s * 0.20, y + s * 0.25, x + s * 0.76, y + s * 0.16, color, 1.0)
    line(slide, x + s * 0.20, y + s * 0.25, x + s * 0.42, y + s * 0.78, color, 1.0)
    line(slide, x + s * 0.76, y + s * 0.16, x + s * 0.42, y + s * 0.78, color, 1.0)
    circle(slide, x + s * 0.07, y + s * 0.12, s * 0.26, color, None)
    circle(slide, x + s * 0.63, y + s * 0.03, s * 0.26, color, None)
    circle(slide, x + s * 0.29, y + s * 0.65, s * 0.26, color, None)


def draw_shield_icon(slide, x, y, s=0.28, color=MINT):
    shield = add_shape(slide, MSO_SHAPE.PENTAGON, x, y, s, s * 1.15, None, color, 1.4)
    shield.rotation = 180
    line(slide, x + s * 0.28, y + s * 0.53, x + s * 0.45, y + s * 0.70, color, 1.1)
    line(slide, x + s * 0.45, y + s * 0.70, x + s * 0.76, y + s * 0.35, color, 1.1)


def draw_lock(slide, x, y, s=0.24, color=AMBER):
    rect(slide, x, y + s * 0.38, s, s * 0.72, None, color, radius=True, line_width=1.3)
    arc = add_shape(slide, MSO_SHAPE.ARC, x + s * 0.18, y, s * 0.64, s * 0.62, None, color, 1.3)
    arc.rotation = 180
    circle(slide, x + s * 0.45, y + s * 0.60, s * 0.11, color, None)


def node_card(slide, x, y, w, h, index, title, subtitle, accent=BLUE, icon=None, fill=SURFACE_2):
    rect(slide, x, y, w, h, fill, C("2B3A56"), radius=True)
    rect(slide, x, y, 0.045, h, accent, None, radius=False)
    if icon == "code":
        draw_code_icon(slide, x + 0.22, y + 0.23, 0.27, accent)
    elif icon == "graph":
        draw_graph_icon(slide, x + 0.20, y + 0.20, 0.31, accent)
    elif icon == "shield":
        draw_shield_icon(slide, x + 0.22, y + 0.19, 0.27, accent)
    elif icon == "lock":
        draw_lock(slide, x + 0.22, y + 0.20, 0.27, accent)
    else:
        circle(slide, x + 0.22, y + 0.24, 0.22, accent, None)
    text(slide, index, x + 0.57, y + 0.17, 0.35, 0.15, 8, accent, True, caps=True)
    text(slide, title, x + 0.57, y + 0.37, w - 0.75, 0.28, 13, TEXT, True)
    text(slide, subtitle, x + 0.57, y + 0.71, w - 0.75, h - 0.79, 8.8, MUTED, line_spacing=1.0)


def mini_step(slide, x, y, w, title, subtitle, accent=BLUE, number=None, h=0.90):
    rect(slide, x, y, w, h, SURFACE_2, C("293851"), radius=True)
    circle(slide, x + 0.13, y + 0.16, 0.29, accent, None)
    text(slide, str(number) if number is not None else "•", x + 0.13, y + 0.207, 0.29, 0.12, 7.7, BG, True, align=PP_ALIGN.CENTER)
    text(slide, title, x + 0.54, y + 0.15, w - 0.66, 0.19, 10.2, TEXT, True)
    text(slide, subtitle, x + 0.54, y + 0.42, w - 0.66, h - 0.48, 8.1, MUTED, line_spacing=1.0)


def dot_matrix(slide, x, y, cols, rows, dx=0.13, dy=0.13, color=DARK_ACCENT):
    for col in range(cols):
        for row in range(rows):
            circle(slide, x + col * dx, y + row * dy, 0.025, color, None)


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    # Quiet halo behind the system mark.
    circle(slide, 8.65, 1.05, 4.65, C("132340"), None, transparency=72)
    circle(slide, 9.20, 1.60, 3.55, C("1B3760"), None, transparency=78)
    dot_matrix(slide, 8.25, 1.22, 16, 16, 0.16, 0.16, C("203452"))
    line(slide, 0.82, 0.54, 12.52, 0.54, C("1C2A42"), 0.7)
    text(slide, "AURA / 00", 0.84, 0.72, 2.0, 0.18, 8.2, CYAN, True, caps=True)
    text(slide, "AURA", 0.82, 1.48, 5.2, 0.74, 53, TEXT, True, font=DISPLAY)
    text(slide, "AI Native Engineering\nOperating Environment", 0.86, 2.48, 6.3, 1.08, 28, TEXT, True, font=DISPLAY, line_spacing=0.93)
    text(slide, "A project-first workspace for understanding systems,\nreasoning over evidence, and governing change.", 0.88, 3.84, 5.3, 0.58, 14, MUTED, line_spacing=1.05)
    rect(slide, 0.88, 5.14, 4.45, 0.04, BLUE, None, radius=False)
    text(slide, "THE UNIT OF WORK IS THE PROJECT — NOT THE PROMPT.", 0.88, 5.38, 5.45, 0.25, 9.3, CYAN, True, caps=True)
    text(slide, "Architecture brief  ·  current working tree  ·  31 July 2026", 0.88, 6.75, 5.6, 0.16, 8.2, FAINT)
    draw_aura_mark(slide, 10.98, 3.18, 1.30, BLUE)
    # Four surrounding system domains.
    for x, y, label, accent in [
        (8.10, 1.48, "PROJECT", CYAN),
        (11.68, 1.67, "KNOWLEDGE", BLUE),
        (8.20, 4.96, "REASONING", PURPLE),
        (11.56, 4.74, "APPROVAL", MINT),
    ]:
        rect(slide, x, y, 1.30, 0.44, SURFACE_2, C("2C4065"), radius=True)
        circle(slide, x + 0.12, y + 0.15, 0.12, accent, None)
        text(slide, label, x + 0.33, y + 0.14, 0.87, 0.14, 7.6, TEXT, True, align=PP_ALIGN.CENTER, caps=True)
    line(slide, 9.35, 1.88, 10.44, 2.53, C("2E4A75"), 1.0)
    line(slide, 11.67, 2.11, 11.60, 2.53, C("2E4A75"), 1.0)
    line(slide, 9.48, 5.18, 10.46, 3.86, C("2E4A75"), 1.0)
    line(slide, 11.55, 4.88, 11.55, 3.87, C("2E4A75"), 1.0)
    line(slide, 0.82, 7.08, 12.52, 7.08, C("1B2639"), 0.7)
    text(slide, "AURA HUB  /  ENGINEERING OPERATING ENVIRONMENT", 0.84, 7.18, 6.0, 0.15, 7.3, FAINT, True, caps=True)
    text(slide, "01", 12.08, 7.15, 0.44, 0.17, 8.5, MUTED, True, align=PP_ALIGN.RIGHT)


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 2, "01 / THE GAP", "The Problem", "A chat window can answer a prompt. It cannot hold a project.")
    text(slide, "Today’s AI coding assistant is optimized for\nconversation — not engineering continuity.", 0.82, 2.03, 5.0, 0.86, 22, TEXT, True, font=DISPLAY, line_spacing=0.95)
    text(slide, "The missing layer is not another model. It is the environment around the model: durable context, evidence, planning, and controlled change.", 0.84, 3.18, 4.72, 0.72, 11.2, MUTED, line_spacing=1.07)
    rect(slide, 0.84, 4.44, 4.5, 1.10, SURFACE_2, C("2A3A58"), radius=True)
    text(slide, "PROMPT  →  ANSWER", 1.08, 4.70, 3.92, 0.24, 12, CYAN, True, align=PP_ALIGN.CENTER, caps=True)
    text(slide, "Useful in the moment. Fragile across the work.", 1.05, 5.06, 4.0, 0.18, 9.2, MUTED, align=PP_ALIGN.CENTER)
    # Limitation cards.
    cards = [
        ("CHAT-SHAPED", "Conversation is the primary unit; the repository is background.", CYAN),
        ("STATELESS", "Context evaporates between tasks, sessions, and engineers.", PURPLE),
        ("PROMPT-LOCAL", "No durable graph of symbols, dependencies, health, or hotspots.", BLUE),
        ("UNPLANNED", "No goal graph, checkpoints, or governed path from intent to change.", AMBER),
    ]
    for i, (label, body, accent) in enumerate(cards):
        x = 6.15 + (i % 2) * 3.18
        y = 2.02 + (i // 2) * 1.78
        rect(slide, x, y, 2.88, 1.38, SURFACE_2, C("293852"), radius=True)
        circle(slide, x + 0.22, y + 0.22, 0.18, accent, None)
        text(slide, label, x + 0.54, y + 0.20, 2.05, 0.17, 8.5, accent, True, caps=True)
        text(slide, body, x + 0.22, y + 0.58, 2.38, 0.55, 10.1, TEXT, True, line_spacing=1.02)
    text(slide, "THE DESIGN QUESTION", 6.16, 5.75, 2.5, 0.15, 8.2, CYAN, True, caps=True)
    text(slide, "How do we make AI useful across the lifecycle of a system —\nwithout surrendering engineering judgment?", 6.15, 6.01, 5.75, 0.48, 13, TEXT, True, font=DISPLAY, line_spacing=0.97)


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 3, "02 / THE VISION", "Make the environment intelligent before making it autonomous.", "AURA changes the unit of work from a prompt to a project.", "NORTH STAR")
    rect(slide, 0.84, 2.03, 3.1, 3.45, SURFACE_2, C("2C3B5C"), radius=True)
    draw_aura_mark(slide, 2.38, 3.33, 0.72, BLUE)
    text(slide, "PROJECT\nAS SYSTEM", 1.13, 4.23, 2.5, 0.60, 19, TEXT, True, font=DISPLAY, align=PP_ALIGN.CENTER, line_spacing=0.93)
    text(slide, "AURA keeps the repository, its structure, its decisions, and its changes in one governed environment.", 1.13, 5.06, 2.5, 0.30, 8.9, MUTED, align=PP_ALIGN.CENTER)
    arrow(slide, 4.05, 3.70, 4.70, 3.70, CYAN, 1.5)
    text(slide, "THREE INVERSIONS", 5.00, 2.03, 3.0, 0.16, 8.2, CYAN, True, caps=True)
    inversions = [
        ("01", "KNOWLEDGE BEFORE\nGENERATION", "The graph and project facts are assembled before the provider is called.", CYAN),
        ("02", "DETERMINISTIC BEFORE\nPROBABILISTIC", "Signals, checks, limits, and arithmetic constrain model judgment.", BLUE),
        ("03", "PROPOSAL BEFORE\nMUTATION", "Analyze, diagnose, and plan freely; write only after explicit human acceptance.", MINT),
    ]
    for i, (number, title, body, accent) in enumerate(inversions):
        y = 2.42 + i * 1.15
        rect(slide, 5.00, y, 6.98, 0.91, SURFACE_2, C("293852"), radius=True)
        circle(slide, 5.22, y + 0.22, 0.44, accent, None)
        text(slide, number, 5.22, y + 0.36, 0.44, 0.12, 7.8, BG, True, align=PP_ALIGN.CENTER)
        text(slide, title, 5.86, y + 0.16, 2.45, 0.45, 10.3, TEXT, True, line_spacing=0.92)
        text(slide, body, 8.52, y + 0.19, 3.13, 0.45, 9.3, MUTED, line_spacing=1.0)
    rect(slide, 5.00, 5.97, 6.98, 0.52, C("152B3A"), C("2C5C64"), radius=True)
    text(slide, "NOT AN AUTONOMOUS EDITOR  ·  A GOVERNED ENGINEERING ENVIRONMENT", 5.24, 6.14, 6.5, 0.16, 8.5, MINT, True, align=PP_ALIGN.CENTER, caps=True)


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 4, "03 / SYSTEM ARCHITECTURE", "A vertical system of understanding, reasoning, and control.", "The desktop is the frame; intelligence is layered around the project.", "ARCHITECTURE")
    # Main vertical spine.
    x, w = 4.05, 5.25
    layers = [
        ("01", "DEVELOPER", "intent · selection · approval", CYAN),
        ("02", "WORKSPACE", "React shell · Monaco · context", BLUE),
        ("03", "KNOWLEDGE FABRIC", "coding index · full-stack graph", PURPLE),
        ("04", "ENGINEERING INTELLIGENCE", "actions · diagnosis · project facts", BLUE),
        ("05", "MISSION CONTROL", "goal graph · DAG · checkpoints", AMBER),
        ("06", "GOVERNANCE", "health · audits · quality gates", MINT),
    ]
    ys = [1.98, 2.72, 3.46, 4.20, 4.94, 5.68]
    for i, ((number, title, sub, accent), y) in enumerate(zip(layers, ys)):
        rect(slide, x, y, w, 0.56, SURFACE_2, C("30405C"), radius=True)
        rect(slide, x, y, 0.06, 0.56, accent, None, radius=False)
        text(slide, number, x + 0.22, y + 0.19, 0.32, 0.13, 7.8, accent, True)
        text(slide, title, x + 0.72, y + 0.13, 2.55, 0.18, 10.8, TEXT, True, caps=True)
        text(slide, sub, x + 3.13, y + 0.18, 1.8, 0.16, 8.4, MUTED, align=PP_ALIGN.RIGHT)
        if i < len(layers) - 1:
            arrow(slide, x + w / 2, y + 0.58, x + w / 2, ys[i + 1] - 0.06, C("3D5686"), 1.1, 0.08)
    # Side rails.
    rect(slide, 0.84, 2.12, 2.50, 1.33, SURFACE, C("27364F"), radius=True)
    text(slide, "LOCAL PROJECT", 1.06, 2.38, 1.8, 0.16, 8.2, CYAN, True, caps=True)
    text(slide, "Files  ·  Git  ·  profiles\n~/.aura JSON stores", 1.06, 2.72, 1.85, 0.42, 10.6, TEXT, True, line_spacing=1.0)
    arrow(slide, 3.38, 2.79, 4.00, 2.79, CYAN, 1.0, 0.08)
    rect(slide, 0.84, 4.23, 2.50, 1.33, SURFACE, C("27364F"), radius=True)
    draw_graph_icon(slide, 1.07, 4.54, 0.33, PURPLE)
    text(slide, "PROJECT TRUTH", 1.55, 4.46, 1.45, 0.16, 8.2, PURPLE, True, caps=True)
    text(slide, "Facts before opinions\nEvidence before action", 1.06, 4.85, 1.85, 0.40, 10.6, TEXT, True, line_spacing=1.0)
    arrow(slide, 3.38, 4.86, 4.00, 4.86, PURPLE, 1.0, 0.08)
    rect(slide, 10.02, 2.12, 2.48, 1.33, SURFACE, C("27364F"), radius=True)
    text(slide, "BYO AI RUNTIME", 10.25, 2.39, 1.85, 0.16, 8.2, BLUE, True, caps=True)
    text(slide, "Provider adapters\nstreaming · model discovery", 10.25, 2.72, 1.84, 0.42, 10.1, TEXT, True, line_spacing=1.0)
    arrow(slide, 10.00, 4.46, 9.38, 4.46, BLUE, 1.0, 0.08)
    rect(slide, 10.02, 4.23, 2.48, 1.33, SURFACE, C("27364F"), radius=True)
    draw_shield_icon(slide, 10.26, 4.53, 0.27, MINT)
    text(slide, "TAURI V2 HOST", 10.70, 4.46, 1.40, 0.16, 8.2, MINT, True, caps=True)
    text(slide, "Thin Rust core\nfilesystem boundary", 10.25, 4.85, 1.85, 0.40, 10.1, TEXT, True, line_spacing=1.0)
    arrow(slide, 10.00, 5.01, 9.38, 5.01, MINT, 1.0, 0.08)
    text(slide, "The architecture is additive: every subsystem reads from the same project context and returns control to the engineer.", 0.84, 6.48, 11.7, 0.25, 10.2, MUTED, italic=True)


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 5, "04 / KNOWLEDGE FABRIC", "The project becomes queryable before it becomes generative.", "Two real engines turn files into searchable context and system-level relationships.", "LIVE TODAY")
    # Coding lane.
    rect(slide, 0.84, 2.02, 5.76, 2.20, SURFACE, C("2B3B58"), radius=True)
    text(slide, "CODING KNOWLEDGE ENGINE", 1.10, 2.26, 3.0, 0.17, 8.7, CYAN, True, caps=True)
    text(slide, "Files → chunks → lexical retrieval → context", 1.10, 2.52, 4.8, 0.24, 12.2, TEXT, True, font=DISPLAY)
    for x, title, sub, accent in [
        (1.11, "SCAN", "filesystem + ignores", CYAN),
        (2.67, "CHUNK", "language-aware", BLUE),
        (4.23, "BM25", "exact / prefix / fuzzy", PURPLE),
    ]:
        rect(slide, x, 3.12, 1.30, 0.61, SURFACE_2, C("30415E"), radius=True)
        circle(slide, x + 0.12, 3.32, 0.16, accent, None)
        text(slide, title, x + 0.34, 3.22, 0.82, 0.15, 8.1, TEXT, True, caps=True)
        text(slide, sub, x + 0.12, 3.48, 1.05, 0.12, 7.3, MUTED, align=PP_ALIGN.CENTER)
        if x < 4.23:
            arrow(slide, x + 1.36, 3.42, x + 1.52, 3.42, C("3D5887"), 1.0, 0.07)
    text(slide, "neighboring chunks  ·  token budget  ·  source refs", 1.10, 3.98, 4.8, 0.15, 8.2, MUTED)
    # Full-stack lane.
    rect(slide, 6.78, 2.02, 5.72, 2.20, SURFACE, C("2B3B58"), radius=True)
    text(slide, "FULL-STACK KNOWLEDGE ENGINE", 7.04, 2.26, 3.1, 0.17, 8.7, PURPLE, True, caps=True)
    text(slide, "Extractors → typed entities → relation graph", 7.04, 2.52, 4.9, 0.24, 12.2, TEXT, True, font=DISPLAY)
    for x, title, sub, accent in [
        (7.05, "EXTRACT", "5 layer extractors", PURPLE),
        (8.61, "LINK", "cross-layer edges", BLUE),
        (10.17, "QUERY", "graph-aware search", CYAN),
    ]:
        rect(slide, x, 3.12, 1.30, 0.61, SURFACE_2, C("30415E"), radius=True)
        circle(slide, x + 0.12, 3.32, 0.16, accent, None)
        text(slide, title, x + 0.34, 3.22, 0.82, 0.15, 8.1, TEXT, True, caps=True)
        text(slide, sub, x + 0.12, 3.48, 1.05, 0.12, 7.3, MUTED, align=PP_ALIGN.CENTER)
        if x < 10.17:
            arrow(slide, x + 1.36, 3.42, x + 1.52, 3.42, C("3D5887"), 1.0, 0.07)
    text(slide, "frontend  ·  backend  ·  database  ·  config  ·  architecture", 7.04, 3.98, 4.95, 0.15, 8.2, MUTED)
    # Shared output.
    arrow(slide, 3.72, 4.33, 3.72, 4.78, CYAN, 1.0, 0.08)
    arrow(slide, 9.62, 4.33, 9.62, 4.78, PURPLE, 1.0, 0.08)
    rect(slide, 1.54, 4.84, 10.30, 1.08, C("12243A"), C("2B5A77"), radius=True)
    draw_graph_icon(slide, 1.88, 5.17, 0.35, CYAN)
    text(slide, "CONTEXT PACKAGE", 2.43, 5.05, 2.0, 0.18, 9.0, CYAN, True, caps=True)
    text(slide, "files  ·  symbols  ·  relations  ·  neighbors  ·  source refs  ·  token budget", 2.43, 5.39, 7.9, 0.16, 10.2, TEXT, True)
    status_pill(slide, "KEYWORD + GRAPH RETRIEVAL IS LIVE", 1.54, 6.22, 3.08, MINT)
    status_pill(slide, "VECTOR / EMBEDDING PATH IS A RESERVED SEAM", 4.84, 6.22, 4.10, AMBER)
    text(slide, "The current engines are deterministic and provider-independent; the model receives assembled context, not the raw repository.", 9.24, 6.19, 3.24, 0.34, 8.4, MUTED, italic=True)


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 6, "05 / AI CODE INTELLIGENCE", "Context is a first-class input to every code action.", "Selection is the starting point; project relationships are what make the response useful.", "IMPLEMENTED")
    steps = [
        ("01", "SELECTION", "Monaco range\nfile + symbol", CYAN),
        ("02", "CONTEXT", "dependencies\ndependents + graph", BLUE),
        ("03", "REASONING", "provider-neutral\nruntime call", PURPLE),
        ("04", "RESPONSE", "explain · diff\nnew file · findings", AMBER),
        ("05", "APPROVAL", "accept / reject\nsave + reindex", MINT),
    ]
    x0 = 0.84
    for i, (n, title, sub, accent) in enumerate(steps):
        x = x0 + i * 2.38
        rect(slide, x, 2.10, 2.02, 1.55, SURFACE_2, C("2C3B57"), radius=True)
        circle(slide, x + 0.20, 2.30, 0.38, accent, None)
        text(slide, n, x + 0.20, 2.43, 0.38, 0.12, 7.8, BG, True, align=PP_ALIGN.CENTER)
        text(slide, title, x + 0.70, 2.27, 1.06, 0.17, 9.2, accent, True, caps=True)
        text(slide, sub, x + 0.20, 2.94, 1.55, 0.42, 10.5, TEXT, True, line_spacing=0.95)
        if i < len(steps) - 1:
            arrow(slide, x + 2.05, 2.87, x + 2.31, 2.87, C("496589"), 1.1, 0.08)
    # Context panel mock.
    rect(slide, 0.84, 4.22, 4.06, 1.62, C("101A2A"), C("294363"), radius=True)
    text(slide, "PROJECT CONTEXT / CURRENT SYMBOL", 1.11, 4.47, 3.35, 0.16, 8.0, CYAN, True, caps=True)
    text(slide, "fetchOrders", 1.11, 4.78, 2.0, 0.24, 15.5, TEXT, True, font=DISPLAY)
    for i, (label, value, accent) in enumerate([("imports", "4", BLUE), ("dependents", "7", PURPLE), ("risk floor", "medium", AMBER)]):
        xx = 1.11 + i * 1.15
        text(slide, label, xx, 5.25, 0.95, 0.13, 7.1, MUTED, caps=True)
        text(slide, value, xx, 5.46, 0.95, 0.18, 10.5, accent, True)
    text(slide, "AURA context is already held client-side before the provider request.", 5.33, 4.36, 6.88, 0.28, 14.5, TEXT, True, font=DISPLAY)
    text(slide, "The action surface is deliberately small: explain, refactor, optimize, diagnose, generate tests, add docs, security review, simplify, convert, rename, or custom.", 5.35, 4.91, 5.85, 0.56, 10.6, MUTED, line_spacing=1.04)
    rect(slide, 5.35, 5.72, 6.82, 0.48, C("162A2D"), C("2A5F5A"), radius=True)
    text(slide, "NO AUTO-WRITE  ·  ACTION → PREVIEW → HUMAN ACCEPTANCE", 5.60, 5.88, 6.30, 0.15, 8.4, MINT, True, align=PP_ALIGN.CENTER, caps=True)


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 7, "06 / ENGINEERING DIAGNOSIS ENGINE", "Diagnosis is an evidence pipeline, not a guess.", "Ten stages narrow uncertainty, compare candidate patches, and stop before an unsafe write.", "IMPLEMENTED")
    stages = [
        ("01", "GATHER", "filesystem · graph\ngit · compiler", CYAN),
        ("02", "CLASSIFY", "deterministic\ncategory + evidence", BLUE),
        ("03", "UNKNOWN?", "honest stop\nno patch", AMBER),
        ("04", "ROOT CAUSE", "AI explanation\nover signals", PURPLE),
        ("05", "A / B / C", "minimal · defensive\nrefactor-adjacent", BLUE),
        ("06", "LIMIT", "diff stats\nexports · layers", CYAN),
        ("07", "SIMULATE", "compiler · refs\ncategory · tests", MINT),
        ("08", "CONFIDENCE", "scores +\nlimiter decision", AMBER),
        ("09", "REVIEW", "AI reviewer\nper candidate", PURPLE),
        ("10", "RECOMMEND", "comparison\nready for human", MINT),
    ]
    for i, (n, title, sub, accent) in enumerate(stages):
        x = 0.84 + (i % 5) * 2.47
        y = 2.06 + (i // 5) * 1.35
        rect(slide, x, y, 2.12, 1.00, SURFACE_2, C("2C3B57"), radius=True)
        circle(slide, x + 0.16, y + 0.17, 0.30, accent, None)
        text(slide, n, x + 0.16, y + 0.28, 0.30, 0.12, 7.0, BG, True, align=PP_ALIGN.CENTER)
        text(slide, title, x + 0.58, y + 0.16, 1.25, 0.15, 8.7, accent, True, caps=True)
        text(slide, sub, x + 0.16, y + 0.54, 1.76, 0.28, 9.0, TEXT, True, line_spacing=0.94)
    # Stage connectors.
    for i in range(4):
        arrow(slide, 2.98 + i * 2.47, 2.56, 3.25 + i * 2.47, 2.56, C("466185"), 0.9, 0.07)
        arrow(slide, 11.46 - i * 2.47, 3.90, 11.19 - i * 2.47, 3.90, C("466185"), 0.9, 0.07)
    arrow(slide, 11.40, 3.07, 11.40, 3.63, C("466185"), 0.9, 0.07)
    rect(slide, 0.84, 5.16, 11.84, 0.76, C("132235"), C("2B5364"), radius=True)
    draw_shield_icon(slide, 1.12, 5.35, 0.25, MINT)
    text(slide, "DECISION SURFACE", 1.56, 5.31, 1.65, 0.15, 8.3, MINT, True, caps=True)
    text(slide, "Candidate compare  ·  Monaco diff  ·  patch limiter  ·  accept / reject  ·  memory record", 3.42, 5.30, 6.90, 0.18, 10.8, TEXT, True)
    text(slide, "Diagnosis runs over SSE stage events and writes nothing until an explicit accept action.", 1.12, 5.65, 9.10, 0.14, 8.4, MUTED)
    status_pill(slide, "NO WRITE DURING ANALYSIS", 10.47, 5.33, 1.90, MINT)
    text(slide, "The current implementation makes “unknown” a valid result — a useful property for an engineering system.", 0.84, 6.29, 11.4, 0.22, 10.0, MUTED, italic=True)


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 8, "07 / MISSION CONTROL", "Intent becomes a governed execution graph.", "Mission Control turns a request into goals, dependencies, waves, checkpoints, and reviewable proposals.", "V3 ENGINE")
    # Planning path.
    text(slide, "PLANNING PIPELINE", 0.84, 2.00, 2.0, 0.16, 8.2, CYAN, True, caps=True)
    plan = [
        ("INTENT", "classify +\nrefine", CYAN),
        ("SIGNAL", "health +\ndebt", BLUE),
        ("GOALS", "goals +\ndeps", PURPLE),
        ("RISK", "risk\nfloor", AMBER),
        ("REVIEW", "quality\nscore", MINT),
    ]
    for i, (title, sub, accent) in enumerate(plan):
        x = 0.84 + i * 1.48
        mini_step(slide, x, 2.35, 1.25, title, sub, accent, number=i + 1, h=0.98)
        if i < len(plan) - 1:
            arrow(slide, x + 1.28, 2.84, x + 1.43, 2.84, C("4D6486"), 1.0, 0.07)
    # Execution graph.
    rect(slide, 0.84, 3.82, 7.80, 2.28, SURFACE, C("2B3A55"), radius=True)
    text(slide, "EXECUTION DAG / DETERMINISTIC WAVES", 1.10, 4.08, 3.8, 0.16, 8.3, AMBER, True, caps=True)
    wave_x = [1.18, 3.15, 5.12, 7.01]
    wave_labels = ["WAVE 01", "WAVE 02", "WAVE 03", "REVIEW"]
    wave_colors = [CYAN, BLUE, PURPLE, MINT]
    for i, (x, label, accent) in enumerate(zip(wave_x, wave_labels, wave_colors)):
        text(slide, label, x, 4.48, 1.20, 0.14, 7.3, accent, True, align=PP_ALIGN.CENTER, caps=True)
        rect(slide, x + 0.10, 4.78, 1.00, 0.42, SURFACE_2, C("30415E"), radius=True)
        circle(slide, x + 0.25, 4.91, 0.14, accent, None)
        text(slide, "task" if i < 3 else "gate", x + 0.45, 4.91, 0.46, 0.12, 7.4, TEXT, True, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow(slide, x + 1.18, 4.99, x + 1.82, 4.99, C("4B648A"), 1.0, 0.08)
    text(slide, "queued  →  waiting  →  running  →  review  →  completed", 1.12, 5.57, 5.7, 0.15, 9.0, MUTED, align=PP_ALIGN.CENTER)
    # Human gates.
    rect(slide, 8.90, 2.35, 3.58, 1.03, C("2D2618"), C("6C5530"), radius=True)
    draw_lock(slide, 9.18, 2.62, 0.26, AMBER)
    text(slide, "PLAN APPROVAL", 9.63, 2.59, 2.1, 0.16, 9.2, AMBER, True, caps=True)
    text(slide, "No task starts before the whole plan is approved.", 9.63, 2.91, 2.45, 0.23, 9.3, TEXT, True)
    rect(slide, 8.90, 3.72, 3.58, 1.03, C("162A2D"), C("2E5C58"), radius=True)
    draw_lock(slide, 9.18, 3.99, 0.26, MINT)
    text(slide, "TASK ACCEPTANCE", 9.63, 3.96, 2.20, 0.16, 9.2, MINT, True, caps=True)
    text(slide, "A proposal is written only after its task is accepted.", 9.63, 4.28, 2.45, 0.23, 9.3, TEXT, True)
    rect(slide, 8.90, 5.13, 3.58, 0.97, SURFACE_2, C("2C3B57"), radius=True)
    text(slide, "NO UNATTENDED EXECUTION", 9.16, 5.42, 3.06, 0.16, 9.4, TEXT, True, align=PP_ALIGN.CENTER, caps=True)
    text(slide, "pause · resume · retry · cancel · replay", 9.16, 5.70, 3.04, 0.14, 8.1, MUTED, align=PP_ALIGN.CENTER)


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 9, "08 / ENGINEERING WORKSPACE", "A surface for the whole loop — not a single editor pane.", "The current desktop brings project navigation, Monaco, context, knowledge, diagnosis, and missions into one frame.", "DESKTOP SURFACE")
    # Workspace chrome.
    rect(slide, 0.84, 2.00, 11.84, 4.46, C("0D1421"), C("2C3B55"), radius=True)
    rect(slide, 0.84, 2.00, 11.84, 0.42, C("111C2C"), C("24334C"), radius=True)
    circle(slide, 1.08, 2.14, 0.08, RED, None)
    circle(slide, 1.24, 2.14, 0.08, AMBER, None)
    circle(slide, 1.40, 2.14, 0.08, MINT, None)
    text(slide, "AURA / CHECKOUT SERVICE", 1.76, 2.13, 2.3, 0.14, 7.9, MUTED, True, caps=True)
    text(slide, "⌘K  COMMAND BAR", 10.80, 2.13, 1.42, 0.14, 7.2, FAINT, True, align=PP_ALIGN.RIGHT, caps=True)
    # Left navigation rail.
    rect(slide, 0.84, 2.42, 1.06, 4.04, C("101927"), None, radius=False)
    draw_aura_mark(slide, 1.37, 2.78, 0.24, BLUE)
    for i, (label, accent) in enumerate([("HOME", FAINT), ("PROJECT", BLUE), ("KNOW", CYAN), ("AI", PURPLE), ("MISSION", AMBER), ("GOV", MINT)]):
        yy = 3.22 + i * 0.48
        if label == "PROJECT":
            rect(slide, 0.99, yy - 0.08, 0.76, 0.31, C("20315A"), None, radius=True)
        circle(slide, 1.12, yy, 0.10, accent, None)
        text(slide, label, 1.31, yy - 0.005, 0.45, 0.12, 6.0, TEXT if label == "PROJECT" else MUTED, True, caps=True)
    # Explorer.
    rect(slide, 1.90, 2.42, 2.10, 4.04, C("101927"), C("24334C"), radius=False)
    text(slide, "EXPLORER", 2.14, 2.72, 1.20, 0.14, 7.8, CYAN, True, caps=True)
    text(slide, "src", 2.18, 3.15, 0.65, 0.14, 9.2, TEXT, True)
    for yy, label, indent, accent in [(3.49, "components", 0.18, MUTED), (3.82, "App.tsx", 0.38, BLUE), (4.15, "workspace", 0.18, MUTED), (4.48, "mission", 0.38, AMBER), (4.81, "pipeline.ts", 0.38, PURPLE), (5.14, "index.ts", 0.38, CYAN)]:
        circle(slide, 2.18 + indent, yy + 0.03, 0.07, accent, None)
        text(slide, label, 2.36 + indent, yy, 1.34 - indent, 0.13, 8.0, TEXT if accent != MUTED else MUTED)
    rect(slide, 2.14, 5.70, 1.54, 0.42, C("16283E"), C("2B4C70"), radius=True)
    text(slide, "+ ASK AURA", 2.32, 5.84, 1.20, 0.13, 7.5, CYAN, True, align=PP_ALIGN.CENTER, caps=True)
    # Editor.
    rect(slide, 4.00, 2.42, 4.64, 4.04, C("0B111C"), C("24334C"), radius=False)
    text(slide, "orders.ts", 4.26, 2.72, 1.2, 0.14, 8.0, TEXT, True)
    text(slide, "CODE  ·  DIFF  ·  DIAGNOSE", 6.33, 2.72, 2.0, 0.14, 7.0, FAINT, True, align=PP_ALIGN.RIGHT, caps=True)
    for i, width in enumerate([2.82, 2.24, 3.08, 1.70, 2.90, 2.46, 3.33, 1.60, 2.68, 2.12, 3.14]):
        yy = 3.18 + i * 0.24
        text(slide, f"{i + 12:02d}", 4.24, yy, 0.28, 0.12, 6.8, FAINT, align=PP_ALIGN.RIGHT)
        rect(slide, 4.76, yy + 0.01, width, 0.055, C("365174") if i in (2, 6) else C("27374D"), None, radius=False, transparency=0)
        if i in (2, 6):
            rect(slide, 4.66, yy - 0.07, 3.72, 0.20, C("152A45"), None, radius=False, transparency=42)
    # Context panel.
    rect(slide, 8.64, 2.42, 4.04, 4.04, C("101927"), C("24334C"), radius=False)
    text(slide, "PROJECT CONTEXT", 8.91, 2.72, 1.85, 0.14, 7.8, CYAN, True, caps=True)
    text(slide, "fetchOrders", 8.91, 3.10, 2.4, 0.22, 14.0, TEXT, True, font=DISPLAY)
    for yy, label, value, accent in [(3.58, "symbol", "function", BLUE), (3.98, "dependencies", "4", PURPLE), (4.38, "dependents", "7", CYAN), (4.78, "architecture", "backend", AMBER)]:
        text(slide, label, 8.91, yy, 1.33, 0.13, 7.1, MUTED, caps=True)
        text(slide, value, 11.09, yy, 1.17, 0.13, 8.8, accent, True, align=PP_ALIGN.RIGHT)
    rect(slide, 8.91, 5.31, 3.37, 0.48, C("172A32"), C("2B5A59"), radius=True)
    text(slide, "OPEN DIAGNOSIS  →", 9.17, 5.48, 2.84, 0.13, 8.1, MINT, True, align=PP_ALIGN.CENTER, caps=True)
    # Bottom mission tray.
    rect(slide, 4.00, 5.96, 4.64, 0.50, C("162137"), C("2B456B"), radius=True)
    text(slide, "MISSION CONTROL", 4.23, 6.13, 1.45, 0.13, 7.6, AMBER, True, caps=True)
    text(slide, "3 tasks  ·  1 review  ·  0 writes", 5.89, 6.13, 2.36, 0.13, 7.8, TEXT, True, align=PP_ALIGN.RIGHT)
    text(slide, "One fixed frame. Many project surfaces. State-driven navigation keeps the environment coherent.", 0.84, 6.70, 11.4, 0.18, 9.5, MUTED, italic=True)


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 10, "09 / SAFETY ARCHITECTURE", "Intelligence can be fast. Mutation must be deliberate.", "AURA separates analysis, proposal, approval, and file writes into inspectable boundaries.", "HUMAN GATED")
    # Flow band.
    flow = [
        ("ANALYZE", "facts + signals", CYAN),
        ("PLAN", "goals + DAG", BLUE),
        ("PROPOSE", "diff + explanation", PURPLE),
        ("WRITE", "accept + reindex", MINT),
    ]
    for i, (title, sub, accent) in enumerate(flow):
        x = 0.86 + i * 2.35
        rect(slide, x, 2.05, 1.88, 0.89, SURFACE_2, C("2C3B57"), radius=True)
        circle(slide, x + 0.17, 2.30, 0.18, accent, None)
        text(slide, title, x + 0.49, 2.20, 1.12, 0.16, 9.4, TEXT, True, caps=True)
        text(slide, sub, x + 0.49, 2.52, 1.15, 0.14, 8.1, MUTED)
        if i < len(flow) - 1:
            arrow(slide, x + 1.92, 2.49, x + 2.26, 2.49, C("456187"), 1.1, 0.08)
    # Gate labels.
    rect(slide, 5.44, 3.21, 2.12, 0.46, C("30291A"), C("755B2E"), radius=True)
    draw_lock(slide, 5.68, 3.33, 0.18, AMBER)
    text(slide, "APPROVE PLAN", 6.01, 3.37, 1.34, 0.13, 8.0, AMBER, True, align=PP_ALIGN.CENTER, caps=True)
    rect(slide, 7.78, 3.21, 2.12, 0.46, C("172D2D"), C("2B625B"), radius=True)
    draw_lock(slide, 8.02, 3.33, 0.18, MINT)
    text(slide, "ACCEPT TASK", 8.35, 3.37, 1.34, 0.13, 8.0, MINT, True, align=PP_ALIGN.CENTER, caps=True)
    line(slide, 6.50, 2.96, 6.50, 3.21, AMBER, 1.0)
    line(slide, 8.84, 2.96, 8.84, 3.21, MINT, 1.0)
    # Safeguard cards.
    safeguards = [
        ("RISK ENGINE", "Deterministic risk floors prevent the model from downgrading a dangerous change.", AMBER),
        ("DIFF + SIMULATION", "A/B/C candidates are compared, limited, and checked before acceptance.", BLUE),
        ("PATH GUARD", "Writes resolve inside the project root; traversal outside the project is rejected.", CYAN),
        ("REINDEX + MEMORY", "Accepted changes reindex the project and record the engineering decision.", MINT),
    ]
    for i, (title, body, accent) in enumerate(safeguards):
        x = 0.86 + (i % 2) * 6.02
        y = 4.13 + (i // 2) * 1.05
        rect(slide, x, y, 5.48, 0.82, SURFACE, C("2B3A55"), radius=True)
        circle(slide, x + 0.23, y + 0.25, 0.18, accent, None)
        text(slide, title, x + 0.57, y + 0.18, 1.55, 0.15, 8.3, accent, True, caps=True)
        text(slide, body, x + 2.25, y + 0.16, 2.96, 0.36, 9.0, TEXT, True, line_spacing=1.0)
    rect(slide, 4.02, 6.30, 5.24, 0.40, C("172B33"), C("2E625A"), radius=True)
    text(slide, "HUMAN IN CONTROL  ·  NO UNATTENDED FILE MUTATION", 4.28, 6.43, 4.72, 0.14, 8.9, MINT, True, align=PP_ALIGN.CENTER, caps=True)


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 11, "10 / ENGINEERING INTELLIGENCE PLATFORM", "One project truth. Multiple intelligence surfaces.", "Knowledge, diagnosis, missions, and workspace are different views over the same engineering context.", "SYSTEM VIEW")
    # Source of truth.
    rect(slide, 0.84, 2.21, 2.35, 3.28, C("12253A"), C("2C5C75"), radius=True)
    draw_graph_icon(slide, 1.72, 2.61, 0.48, CYAN)
    text(slide, "KNOWLEDGE\nFABRIC", 1.11, 3.33, 1.80, 0.53, 20, TEXT, True, font=DISPLAY, align=PP_ALIGN.CENTER, line_spacing=0.92)
    text(slide, "facts · entities\nrelations · context", 1.14, 4.27, 1.75, 0.38, 10.0, MUTED, align=PP_ALIGN.CENTER, line_spacing=1.0)
    # Fan-out.
    for yy in (2.48, 3.48, 4.48):
        arrow(slide, 3.27, 3.85, 3.91, yy + 0.38, C("3D6082"), 1.1, 0.07)
    surfaces = [
        ("AI CODE INTELLIGENCE", "selection → context → action", BLUE),
        ("DIAGNOSIS ENGINE", "signals → candidates → review", PURPLE),
        ("MISSION CONTROL", "intent → graph → execution", AMBER),
    ]
    for i, (title, sub, accent) in enumerate(surfaces):
        y = 2.21 + i * 1.00
        rect(slide, 4.00, y, 3.70, 0.76, SURFACE_2, C("2C3B57"), radius=True)
        rect(slide, 4.00, y, 0.06, 0.76, accent, None, radius=False)
        text(slide, title, 4.29, y + 0.15, 2.92, 0.15, 9.0, TEXT, True, caps=True)
        text(slide, sub, 4.29, y + 0.43, 2.98, 0.14, 8.7, MUTED)
        arrow(slide, 7.75, y + 0.38, 8.43, y + 0.38, accent, 1.0, 0.07)
    # Shared output.
    rect(slide, 8.50, 2.21, 3.98, 3.28, C("131F2D"), C("2F465D"), radius=True)
    text(slide, "SHARED GOVERNED LOOP", 8.80, 2.52, 2.8, 0.16, 8.3, MINT, True, caps=True)
    loop = [("DIFF", BLUE), ("APPROVAL", AMBER), ("EXECUTION", PURPLE), ("REINDEX", CYAN), ("MEMORY", MINT)]
    for i, (label, accent) in enumerate(loop):
        y = 2.94 + i * 0.40
        circle(slide, 8.84, y + 0.02, 0.12, accent, None)
        text(slide, label, 9.10, y, 1.28, 0.14, 8.6, TEXT, True, caps=True)
        if i < len(loop) - 1:
            line(slide, 8.90, y + 0.15, 8.90, y + 0.40, C("39506A"), 0.8)
    text(slide, "The provider is a verbalizer over AURA context — not the source of truth.", 9.10, 5.02, 2.95, 0.24, 9.3, MUTED, italic=True)
    rect(slide, 0.84, 5.96, 11.64, 0.56, C("162A39"), C("2C5B73"), radius=True)
    text(slide, "ONE GRAPH  ·  DIFFERENT CONTEXT PACKS  ·  EXPLICIT HUMAN CONTROL", 1.16, 6.15, 11.02, 0.16, 9.2, CYAN, True, align=PP_ALIGN.CENTER, caps=True)


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 12, "11 / CURRENT STATUS", "What is real in the current working tree.", "Implemented surfaces only. No roadmap concepts are included on this slide.", "IMPLEMENTED")
    items = [
        ("DESKTOP SHELL", "React/Vite · Zustand · command palette · theme tokens · Tauri v2 host", BLUE),
        ("PROJECT LIFECYCLE", "add / open / profile · session restore · project-scoped workspace", CYAN),
        ("CODING KNOWLEDGE", "recursive + incremental index · BM25 search · token-budgeted context", CYAN),
        ("FULL-STACK GRAPH", "five extractors · typed entities / relations · persistent graph-aware search", PURPLE),
        ("BYO AI RUNTIME", "provider adapters · streaming · model discovery · local HTTP/SSE service", BLUE),
        ("AI CODE WORKSPACE", "context panel · grounded actions · diff / new-file / findings", PURPLE),
        ("DIAGNOSIS ENGINE", "real signals · A/B/C candidates · simulation · review · accept/reject", AMBER),
        ("MISSION CONTROL", "goal graph · DAG waves · checkpoints · replay · global dashboard", AMBER),
        ("WORKFLOWS + MEMORY", "visual workflow store/run path · basic engineering decision/event records", MINT),
    ]
    for i, (title, body, accent) in enumerate(items):
        col = i % 3
        row = i // 3
        x = 0.84 + col * 4.02
        y = 2.05 + row * 1.17
        rect(slide, x, y, 3.70, 0.89, SURFACE_2, C("2A3954"), radius=True)
        circle(slide, x + 0.20, y + 0.20, 0.17, accent, None)
        text(slide, title, x + 0.52, y + 0.16, 2.85, 0.15, 8.3, accent, True, caps=True)
        text(slide, body, x + 0.20, y + 0.46, 3.20, 0.27, 9.2, TEXT, True, line_spacing=1.0)
    rect(slide, 0.84, 5.89, 11.78, 0.56, C("172A2D"), C("2D5A59"), radius=True)
    text(slide, "CURRENT TRUTH  ·  REAL FILES  ·  REAL INDEXES  ·  REAL SSE EVENTS  ·  HUMAN-GATED WRITES", 1.10, 6.08, 11.24, 0.16, 9.0, MINT, True, align=PP_ALIGN.CENTER, caps=True)


def slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 13, "12 / FUTURE ROADMAP", "Build the learning loop — without weakening the gates.", "The next horizon is separated from today’s implemented behavior by design.", "FUTURE / IN PROGRESS")
    roadmap = [
        ("IN PROGRESS", "ENGINEERING MEMORY", "Durable decision and event records are integrated; pattern, experience, graph, and insight depth are still expanding.", MINT),
        ("IN PROGRESS", "GOVERNANCE", "Backend modules and a dashboard surface exist; deeper product integration and enforcement are next.", MINT),
        ("FUTURE", "PREDICTION", "Forecast hotspots, risk, and likely failure modes from accumulated engineering evidence.", BLUE),
        ("FUTURE", "LEARNING", "Close the loop from accepted, rejected, and reverted outcomes — with provenance.", PURPLE),
        ("FUTURE", "DIGITAL TWIN", "Continuously model system state, dependencies, health, and change impact.", CYAN),
        ("FUTURE", "AUTONOMOUS ENGINEERING", "Bounded execution policies may grow; explicit control remains a product principle.", AMBER),
    ]
    for i, (status, title, body, accent) in enumerate(roadmap):
        col = i % 2
        row = i // 2
        x = 0.84 + col * 6.02
        y = 2.04 + row * 1.32
        rect(slide, x, y, 5.48, 1.04, SURFACE_2, C("2A3954"), radius=True)
        status_pill(slide, status, x + 0.24, y + 0.18, 1.20 if status == "FUTURE" else 1.56, accent)
        text(slide, title, x + 1.74 if status == "FUTURE" else x + 2.10, y + 0.20, 2.80, 0.15, 9.0, TEXT, True, caps=True)
        text(slide, body, x + 0.24, y + 0.57, 4.86, 0.30, 9.3, MUTED, line_spacing=1.0)
    rect(slide, 0.84, 6.20, 11.78, 0.42, C("2A2116"), C("6A512D"), radius=True)
    text(slide, "NORTH STAR  ·  MORE CAPABLE PROPOSALS  ·  NEVER LESS EXPLICIT CONTROL", 1.10, 6.33, 11.24, 0.14, 8.8, AMBER, True, align=PP_ALIGN.CENTER, caps=True)


def slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 14, "13 / TECHNOLOGY STACK", "A local-first stack with provider-neutral intelligence.", "The architecture keeps the desktop responsive, the AI runtime replaceable, and the project data close to the engineer.", "STACK")
    rows = [
        ("EXPERIENCE", ["React", "TypeScript", "Vite", "Monaco", "@aura/ui"], BLUE),
        ("DESKTOP + SERVICE", ["Tauri v2", "thin Rust", "Node", "localhost HTTP/SSE", "Zustand"], CYAN),
        ("KNOWLEDGE", ["Coding Engine", "FullStack Engine", "BM25", "persistent graph", "graphify"], PURPLE),
        ("AI RUNTIME", ["Anthropic", "OpenAI", "Gemini", "Groq", "Mistral", "NVIDIA", "OpenRouter", "Kimi"], AMBER),
        ("ENGINEERING SYSTEMS", ["Git", "JSON stores", "Diagnosis", "Mission Control", "Governance", "Architecture docs"], MINT),
    ]
    for i, (label, values, accent) in enumerate(rows):
        y = 2.05 + i * 0.87
        text(slide, label, 0.84, y + 0.20, 1.72, 0.16, 8.2, accent, True, caps=True)
        line(slide, 2.58, y + 0.40, 2.92, y + 0.40, C("30415E"), 1.0)
        if len(values) > 6:
            item_step, item_width, item_size = 1.07, 0.96, 7.7
        elif len(values) == 6:
            item_step, item_width, item_size = 1.45, 1.30, 8.0
        else:
            item_step, item_width, item_size = 1.60, 1.43, 8.4
        for j, value in enumerate(values):
            x = 3.08 + j * item_step
            rect(slide, x, y + 0.11, item_width, 0.56, SURFACE_2, C("2B3A55"), radius=True)
            text(slide, value, x + 0.05, y + 0.30, item_width - 0.10, 0.13, item_size, TEXT, True, align=PP_ALIGN.CENTER)
    rect(slide, 0.84, 6.57, 11.78, 0.30, C("182236"), C("2D4568"), radius=True)
    text(slide, "CURRENT RETRIEVAL IS KEYWORD + GRAPH; EMBEDDINGS / VECTOR SEARCH REMAIN A REPLACEABLE SEAM.", 1.06, 6.66, 11.34, 0.12, 7.9, MUTED, True, align=PP_ALIGN.CENTER, caps=True)


def slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base(slide, 15, "14 / LIVE DEMONSTRATION FLOW", "From an open project to an approved change.", "A single path makes the architecture tangible: each step exposes a real artifact or gate.", "DEMO PATH")
    flow = [
        ("01", "OPEN PROJECT", "profile +\nindex status", CYAN),
        ("02", "KNOWLEDGE FABRIC", "files +\ngraph context", BLUE),
        ("03", "ASK AURA", "grounded\nconversation", PURPLE),
        ("04", "DIAGNOSIS", "signals +\nA/B/C patches", AMBER),
        ("05", "MISSION", "goal graph +\nexecution DAG", AMBER),
        ("06", "REVIEW DIFF", "proposal +\nlimiter", BLUE),
        ("07", "APPROVE", "human\nacceptance", MINT),
        ("08", "REINDEX", "memory +\nnew context", CYAN),
    ]
    for i, (n, title, sub, accent) in enumerate(flow):
        row = 0 if i < 4 else 1
        col = i if i < 4 else i - 4
        x = 0.84 + col * 3.00
        y = 2.16 + row * 1.95
        rect(slide, x, y, 2.52, 1.23, SURFACE_2, C("2C3B57"), radius=True)
        circle(slide, x + 0.20, y + 0.20, 0.38, accent, None)
        text(slide, n, x + 0.20, y + 0.33, 0.38, 0.12, 7.8, BG, True, align=PP_ALIGN.CENTER)
        text(slide, title, x + 0.75, y + 0.20, 1.55, 0.15, 8.9, accent, True, caps=True)
        text(slide, sub, x + 0.20, y + 0.66, 1.88, 0.32, 11.1, TEXT, True, line_spacing=0.94)
        if col < 3:
            arrow(slide, x + 2.55, y + 0.62, x + 2.88, y + 0.62, C("456187"), 1.0, 0.08)
    arrow(slide, 10.10, 3.40, 10.10, 4.00, C("456187"), 1.0, 0.08)
    arrow(slide, 3.36, 5.34, 3.36, 5.92, C("456187"), 1.0, 0.08)
    rect(slide, 0.84, 6.10, 11.78, 0.48, C("142A37"), C("2C5D6D"), radius=True)
    text(slide, "REAL FILES  ·  SSE STAGES  ·  SOURCE CONTEXT  ·  HUMAN CLICK BEFORE WRITE", 1.12, 6.26, 11.24, 0.14, 8.8, CYAN, True, align=PP_ALIGN.CENTER, caps=True)


def slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    circle(slide, 8.42, 0.92, 4.78, C("102441"), None, transparency=70)
    circle(slide, 9.09, 1.59, 3.42, C("1A3761"), None, transparency=78)
    dot_matrix(slide, 8.24, 1.18, 17, 17, 0.15, 0.15, C("1E3354"))
    line(slide, 0.82, 0.54, 12.52, 0.54, C("1C2A42"), 0.7)
    text(slide, "AURA / CLOSING", 0.84, 0.72, 2.0, 0.18, 8.2, CYAN, True, caps=True)
    text(slide, "The next engineering\nenvironment is a system.", 0.84, 1.62, 6.0, 1.15, 31, TEXT, True, font=DISPLAY, line_spacing=0.92)
    text(slide, "It understands the project, reasons over evidence,\nand waits for the engineer.", 0.88, 3.25, 5.4, 0.58, 16, MUTED, line_spacing=1.03)
    draw_aura_mark(slide, 10.90, 3.10, 1.36, BLUE)
    for x, title, body, accent in [
        (0.88, "SHARED TRUTH", "The repository becomes queryable context.", CYAN),
        (3.60, "SAFER CHANGE", "Every mutation has a visible gate.", MINT),
        (6.32, "DURABLE CONTEXT", "Decisions can outlive the chat.", PURPLE),
    ]:
        rect(slide, x, 5.28, 2.42, 0.74, SURFACE_2, C("2B3A55"), radius=True)
        circle(slide, x + 0.18, 5.52, 0.15, accent, None)
        text(slide, title, x + 0.46, 5.43, 1.72, 0.14, 8.0, accent, True, caps=True)
        text(slide, body, x + 0.18, 5.74, 1.98, 0.15, 8.5, TEXT, True)
    rect(slide, 9.34, 5.43, 3.15, 0.50, C("172A37"), C("2E5D6D"), radius=True)
    text(slide, "AURA HUB  ·  AI NATIVE ENGINEERING OPERATING ENVIRONMENT", 9.53, 5.60, 2.78, 0.15, 7.7, CYAN, True, align=PP_ALIGN.CENTER, caps=True)
    line(slide, 0.82, 7.08, 12.52, 7.08, C("1B2639"), 0.7)
    text(slide, "AURA HUB  /  ENGINEERING OPERATING ENVIRONMENT  /  REPO-VERIFIED", 0.84, 7.18, 7.6, 0.16, 7.3, FAINT, True, caps=True)
    text(slide, "16", 12.08, 7.15, 0.44, 0.17, 8.5, MUTED, True, align=PP_ALIGN.RIGHT)


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "AURA — AI Native Engineering Operating Environment"
    prs.core_properties.subject = "Repository-grounded architecture presentation"
    prs.core_properties.author = "AURA Hub"
    prs.core_properties.keywords = "AURA, engineering intelligence, knowledge fabric, mission control"
    for maker in (
        slide_1,
        slide_2,
        slide_3,
        slide_4,
        slide_5,
        slide_6,
        slide_7,
        slide_8,
        slide_9,
        slide_10,
        slide_11,
        slide_12,
        slide_13,
        slide_14,
        slide_15,
        slide_16,
    ):
        maker(prs)
    prs.save(OUT)
    print(f"wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
